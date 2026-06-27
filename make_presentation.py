"""
Generates AIJobAnalyzer.pptx — a fully styled PowerPoint presentation.
Run: python make_presentation.py
"""

import io, math, textwrap
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ─── Palette ───────────────────────────────────────────────────────────────
INDIGO    = RGBColor(0x4F, 0x46, 0xE5)   # indigo-600
INDIGO_L  = RGBColor(0xE0, 0xE7, 0xFF)   # indigo-100
INDIGO_D  = RGBColor(0x38, 0x2F, 0xBD)   # indigo-700
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_700 = RGBColor(0x33, 0x4A, 0x5E)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN     = RGBColor(0x16, 0xA3, 0x4A)
RED       = RGBColor(0xDC, 0x26, 0x26)
AMBER     = RGBColor(0xD9, 0x77, 0x06)
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)
EMERALD   = RGBColor(0x05, 0x96, 0x69)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

# ─── Helpers ───────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def rgb_to_hex(r, g, b):
    return "#{:02X}{:02X}{:02X}".format(r, g, b)


def add_rect(slide, x, y, w, h, fill: RGBColor, alpha=None, radius=0):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        x, y, w, h
    )
    # Use rounded rectangle when radius requested
    if radius:
        shape.adjustments[0] = radius / 100
    fill_obj = shape.fill
    fill_obj.solid()
    fill_obj.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, x, y, w, h, fill: RGBColor, radius_emu=Inches(0.15)):
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.oxml.ns import qn
    sp = slide.shapes.add_shape(
        pptx.oxml.ns.qn if False else 5,  # 5 = msoShapeRoundedRectangle
        x, y, w, h
    )
    sp = slide.shapes.add_shape(5, x, y, w, h)  # 5 = Rounded Rectangle
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.adjustments[0] = 0.08   # corner radius ratio
    return sp


def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, color=SLATE_800,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False,
                 font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_multiline_text(slide, lines, x, y, w, h,
                       font_size=16, color=SLATE_700, line_space=1.2,
                       bold=False, font_name="Calibri"):
    """lines: list of strings."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after  = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.name  = font_name
    return txb


def img_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    return buf


def add_fig(slide, fig, x, y, w, h):
    buf = img_to_stream(fig)
    slide.shapes.add_picture(buf, x, y, w, h)
    plt.close(fig)


def pil_to_stream(img: Image.Image):
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ─── Background helpers ────────────────────────────────────────────────────

def dark_bg(slide):
    add_rect(slide, 0, 0, W, H, SLATE_900)


def light_bg(slide):
    add_rect(slide, 0, 0, W, H, SLATE_100)


def white_bg(slide):
    add_rect(slide, 0, 0, W, H, WHITE)


def accent_strip(slide, height=Inches(0.07)):
    add_rect(slide, 0, 0, W, height, INDIGO)


def slide_header(slide, title, subtitle=None, dark=False):
    tc = WHITE if dark else SLATE_900
    sc = RGBColor(0xA5, 0xB4, 0xFC) if dark else SLATE_500
    add_text_box(slide, title, Inches(0.6), Inches(0.25), Inches(12), Inches(0.7),
                 font_size=32, bold=True, color=tc, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, Inches(0.6), Inches(0.95), Inches(12), Inches(0.4),
                     font_size=17, color=sc, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

# ── Slide 1 · Title ────────────────────────────────────────────────────────
def slide_title(prs):
    sl = blank_slide(prs)
    dark_bg(sl)

    # Gradient accent bar left
    add_rect(sl, 0, 0, Inches(0.18), H, INDIGO)

    # Big heading
    add_text_box(sl, "AI Job Analyzer",
                 Inches(0.5), Inches(1.6), Inches(8.5), Inches(1.6),
                 font_size=60, bold=True, color=WHITE, font_name="Calibri")

    add_text_box(sl, "Your AI-powered career co-pilot",
                 Inches(0.5), Inches(3.15), Inches(8), Inches(0.6),
                 font_size=24, color=RGBColor(0xA5, 0xB4, 0xFC), font_name="Calibri")

    add_text_box(sl, "Resume Analysis  ·  Job Tracking  ·  Job Matching  ·  Cover Letter Generation",
                 Inches(0.5), Inches(3.85), Inches(10), Inches(0.5),
                 font_size=15, color=SLATE_500, font_name="Calibri")

    # Right panel — feature icon grid
    fig, ax = plt.subplots(figsize=(4.5, 5.5))
    fig.patch.set_facecolor("#0F172A")
    ax.set_facecolor("#0F172A")
    ax.set_xlim(0, 2); ax.set_ylim(0, 2); ax.axis("off")

    icons = [
        (0.35, 1.55, "#4F46E5", "Resume\nAnalyzer"),
        (1.15, 1.55, "#2563EB", "Job\nTracker"),
        (0.35, 0.65, "#7C3AED", "Job\nMatcher"),
        (1.15, 0.65, "#059669", "Cover\nLetter"),
    ]
    for cx, cy, col, lbl in icons:
        circ = plt.Circle((cx, cy), 0.28, color=col, zorder=3)
        ax.add_patch(circ)
        ax.text(cx, cy - 0.47, lbl, ha="center", va="top",
                fontsize=9, color="white", fontweight="bold",
                multialignment="center")

    ax.text(1.0, 1.97, "4 AI-Powered Tools", ha="center", va="top",
            fontsize=11, color="#A5B4FC", fontweight="bold")
    add_fig(sl, fig, Inches(8.8), Inches(1.2), Inches(4.0), Inches(5.5))

    # Tech pills at bottom
    pills = ["Next.js 14", "FastAPI", "Claude AI", "Tailwind CSS", "SQLite", "Vercel · Render"]
    for i, p in enumerate(pills):
        add_rounded_rect(sl, Inches(0.5 + i * 2.0), Inches(6.55), Inches(1.75), Inches(0.55),
                         INDIGO_D)
        add_text_box(sl, p, Inches(0.5 + i * 2.0), Inches(6.6), Inches(1.75), Inches(0.45),
                     font_size=11, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

    return sl


# ── Slide 2 · The Problem ──────────────────────────────────────────────────
def slide_problem(prs):
    sl = blank_slide(prs)
    white_bg(sl)
    accent_strip(sl)
    slide_header(sl, "The Problem", "Job hunting is painful — and the odds are stacked against you")

    problems = [
        ("#DC2626", "ATS Rejection",
         "Up to 75% of resumes are rejected by automated\ntracking systems before a human ever reads them."),
        ("#D97706", "Zero Visibility",
         "Applicants have no idea how well they fit a role.\nApplying blind wastes months of effort."),
        ("#7C3AED", "Cover Letter Fatigue",
         "Writing a tailored cover letter for every application\ntakes 30–60 minutes per job."),
        ("#2563EB", "Tracking Chaos",
         "Managing dozens of applications in spreadsheets\nis error-prone and demoralizing."),
    ]

    for i, (col, title, desc) in enumerate(problems):
        row, col_n = divmod(i, 2)
        x = Inches(0.5 + col_n * 6.3)
        y = Inches(1.55 + row * 2.65)

        # Card
        card = sl.shapes.add_shape(5, x, y, Inches(5.9), Inches(2.35))
        card.fill.solid(); card.fill.fore_color.rgb = SLATE_100
        card.line.color.rgb = RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)])
        card.line.width = Pt(2)
        card.adjustments[0] = 0.05

        # Accent dot
        dot = sl.shapes.add_shape(9, x + Inches(0.2), y + Inches(0.18),
                                  Inches(0.32), Inches(0.32))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)])
        dot.line.fill.background()

        add_text_box(sl, title, x + Inches(0.6), y + Inches(0.12), Inches(5), Inches(0.42),
                     font_size=16, bold=True, color=SLATE_900)
        add_text_box(sl, desc, x + Inches(0.2), y + Inches(0.65), Inches(5.5), Inches(1.5),
                     font_size=13, color=SLATE_700, wrap=True)

    return sl


# ── Slide 3 · Tech Stack ───────────────────────────────────────────────────
def slide_tech_stack(prs):
    sl = blank_slide(prs)
    white_bg(sl)
    accent_strip(sl)
    slide_header(sl, "Tech Stack", "Modern, production-ready tooling — deployed on free tiers")

    # Matplotlib diagram
    fig, ax = plt.subplots(figsize=(12, 4.8))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.set_xlim(0, 12); ax.set_ylim(0, 5); ax.axis("off")

    layers = [
        (1.0,  3.5, 3.0, "#EEF2FF", "#4F46E5", "FRONTEND",
         ["Next.js 14 (App Router)", "TypeScript", "Tailwind CSS", "Lucide Icons"]),
        (4.5,  3.5, 3.0, "#F5F3FF", "#7C3AED", "BACKEND",
         ["FastAPI (Python)", "SQLAlchemy ORM", "SQLite Database", "Uvicorn ASGI"]),
        (8.0,  3.5, 3.0, "#ECFDF5", "#059669", "AI LAYER",
         ["Claude claude-sonnet-4-6", "Structured JSON output", "Resume parsing", "NLP analysis"]),
    ]

    for lx, ly, lw, bg, col, label, items in layers:
        rect = FancyBboxPatch((lx, ly - len(items)*0.52 - 0.1), lw, len(items)*0.52 + 0.72,
                              boxstyle="round,pad=0.1", linewidth=2,
                              edgecolor=col, facecolor=bg)
        ax.add_patch(rect)
        ax.text(lx + lw/2, ly + 0.2, label, ha="center", va="bottom",
                fontsize=12, fontweight="bold", color=col)
        for j, item in enumerate(items):
            ax.text(lx + 0.25, ly - j*0.52 - 0.05, f"• {item}", va="top",
                    fontsize=9.5, color="#334155")

    # Arrows between layers
    for ax_x in [4.0, 7.5]:
        ax.annotate("", xy=(ax_x + 0.55, 2.5), xytext=(ax_x, 2.5),
                    arrowprops=dict(arrowstyle="->", color="#94A3B8", lw=2))

    # Deploy row
    deploys = [
        (1.0,  0.1, 2.7, "#EEF2FF", "#4F46E5", "Vercel", "Frontend deploy"),
        (4.5,  0.1, 2.7, "#F5F3FF", "#7C3AED", "Render", "Backend deploy"),
        (8.0,  0.1, 2.7, "#FFF7ED", "#EA580C", "Anthropic", "Claude API"),
    ]
    for lx, ly, lw, bg, col, name, sub in deploys:
        rect = FancyBboxPatch((lx, ly), lw, 0.9,
                              boxstyle="round,pad=0.08", linewidth=1.5,
                              edgecolor=col, facecolor=bg)
        ax.add_patch(rect)
        ax.text(lx + lw/2, ly + 0.62, name, ha="center", va="bottom",
                fontsize=11, fontweight="bold", color=col)
        ax.text(lx + lw/2, ly + 0.1, sub, ha="center", va="bottom",
                fontsize=8.5, color="#64748B")

    ax.text(6, 4.85, "Architecture Stack", ha="center", va="top",
            fontsize=14, fontweight="bold", color="#0F172A")

    add_fig(sl, fig, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.8))
    return sl


# ── Slide 4 · Architecture ─────────────────────────────────────────────────
def slide_architecture(prs):
    sl = blank_slide(prs)
    dark_bg(sl)
    add_rect(sl, 0, 0, Inches(0.18), H, INDIGO)
    slide_header(sl, "System Architecture", "How the pieces connect end-to-end", dark=True)

    fig, ax = plt.subplots(figsize=(12.5, 5.2))
    fig.patch.set_facecolor("#0F172A")
    ax.set_facecolor("#0F172A")
    ax.set_xlim(0, 13); ax.set_ylim(0, 5.5); ax.axis("off")

    def box(lx, ly, lw, lh, col, label, items=None, text_col="white"):
        rect = FancyBboxPatch((lx, ly), lw, lh, boxstyle="round,pad=0.12",
                              linewidth=2, edgecolor=col, facecolor=col+"22" if len(col)==7 else col)
        # convert hex shorthand
        ec = col; fc_alpha = col + "33" if "#" in col else col
        rect2 = FancyBboxPatch((lx, ly), lw, lh, boxstyle="round,pad=0.12",
                               linewidth=2, edgecolor=ec,
                               facecolor=(*[int(ec[i:i+2], 16)/255 for i in (1,3,5)], 0.18))
        ax.add_patch(rect2)
        ax.text(lx + lw/2, ly + lh - 0.28, label, ha="center", va="top",
                fontsize=11, fontweight="bold", color=ec)
        if items:
            for j, itm in enumerate(items):
                ax.text(lx + 0.2, ly + lh - 0.65 - j*0.48, f"• {itm}",
                        fontsize=8.5, color="#CBD5E1", va="top")

    def arrow(x1, y1, x2, y2, label=""):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="->", color="#64748B", lw=1.8))
        if label:
            mx, my = (x1+x2)/2, (y1+y2)/2
            ax.text(mx, my + 0.15, label, ha="center", fontsize=7.5, color="#94A3B8")

    # Browser
    box(0.3, 3.5, 3.2, 1.7, "#4F46E5", "BROWSER / Next.js",
        ["Dashboard", "Resume Analyzer", "Job Tracker", "Matcher · Cover Letter"])
    # API proxy
    box(4.0, 3.8, 2.2, 1.1, "#2563EB", "/api proxy", ["App Router catch-all", "CORS handled"])
    # FastAPI
    box(7.0, 2.8, 3.2, 2.4, "#7C3AED", "FastAPI Backend",
        ["routers/resume.py", "routers/jobs.py", "routers/match.py", "routers/cover_letter.py"])
    # Claude
    box(7.0, 0.35, 3.2, 1.8, "#059669", "Claude claude-sonnet-4-6 (Anthropic)",
        ["Structured JSON", "Resume analysis", "Match scoring", "Letter generation"])
    # SQLite
    box(11.0, 2.8, 1.8, 2.4, "#D97706", "SQLite DB",
        ["jobs table", "Persist apps"])

    # localStorage
    box(0.3, 1.2, 3.2, 1.6, "#0F172A", "localStorage",
        ["aijob_resume_text", "aijob_jd_text"])
    # Override border
    rect_ls = FancyBboxPatch((0.3, 1.2), 3.2, 1.6, boxstyle="round,pad=0.12",
                             linewidth=1.5, edgecolor="#64748B",
                             facecolor=(1,1,1,0.04), linestyle="dashed")
    ax.add_patch(rect_ls)
    ax.text(1.9, 2.65, "localStorage", ha="center", va="top",
            fontsize=10, fontweight="bold", color="#94A3B8")
    ax.text(1.9, 2.2, "aijob_resume_text", fontsize=8, color="#64748B", ha="center")
    ax.text(1.9, 1.78, "aijob_jd_text", fontsize=8, color="#64748B", ha="center")

    # Arrows
    arrow(3.5, 4.35, 4.0, 4.35, "fetch")
    arrow(6.2, 4.35, 7.0, 4.35, "REST")
    arrow(8.6, 2.8, 8.6, 2.15, "Claude SDK")
    arrow(10.2, 3.8, 11.0, 3.8, "SQL")
    # localStorage loop
    ax.annotate("", xy=(1.9, 2.8), xytext=(1.9, 3.5),
                arrowprops=dict(arrowstyle="<->", color="#64748B", lw=1.5, linestyle="dashed"))
    ax.text(2.3, 3.15, "persist\nresume", fontsize=7.5, color="#94A3B8", ha="left")

    add_fig(sl, fig, Inches(0.2), Inches(1.35), Inches(12.9), Inches(5.9))
    return sl


# ── Slide 5 · Dashboard ────────────────────────────────────────────────────
def slide_dashboard(prs):
    sl = blank_slide(prs)
    white_bg(sl)
    accent_strip(sl)
    slide_header(sl, "Feature 1 · Dashboard",
                 "A command center for the entire job search")

    # Left: feature bullets
    features = [
        ("Stats Cards", "Total, Applied, Interviewing, Offers, Rejected — live counts from the API"),
        ("Recent Applications", "Latest 5 entries with status badges, link to full tracker"),
        ("AI Activity Log", "Timestamped history of every Claude action, persisted in localStorage"),
        ("Quick Actions", "One-click shortcuts to all four AI tools"),
        ("Pro Tip Card", "Contextual advice surfaced on the dashboard"),
    ]
    for i, (title, desc) in enumerate(features):
        y = Inches(1.5 + i * 1.04)
        dot = sl.shapes.add_shape(9, Inches(0.5), y + Inches(0.08),
                                  Inches(0.22), Inches(0.22))
        dot.fill.solid(); dot.fill.fore_color.rgb = INDIGO
        dot.line.fill.background()
        add_text_box(sl, title, Inches(0.85), y, Inches(5.5), Inches(0.35),
                     font_size=14, bold=True, color=INDIGO)
        add_text_box(sl, desc, Inches(0.85), y + Inches(0.33), Inches(5.5), Inches(0.55),
                     font_size=12, color=SLATE_700)

    # Right: dashboard mockup
    fig, ax = plt.subplots(figsize=(5.8, 5.5))
    fig.patch.set_facecolor("#F8FAFC")
    ax.set_facecolor("#F8FAFC")
    ax.set_xlim(0, 6); ax.set_ylim(0, 5.8); ax.axis("off")

    # Stat cards row
    stats = [("12", "Total", "#4F46E5"), ("5", "Applied", "#2563EB"),
             ("3", "Interview", "#7C3AED"), ("1", "Offer", "#16A34A"), ("3", "Rejected", "#DC2626")]
    for i, (val, lbl, col) in enumerate(stats):
        cx = 0.1 + i * 1.17
        rect = FancyBboxPatch((cx, 4.9), 1.05, 0.75, boxstyle="round,pad=0.06",
                              linewidth=1, edgecolor="#E2E8F0", facecolor="white")
        ax.add_patch(rect)
        ax.text(cx+0.52, 5.42, val, ha="center", va="center",
                fontsize=14, fontweight="bold", color=col)
        ax.text(cx+0.52, 5.07, lbl, ha="center", va="center",
                fontsize=6.5, color="#94A3B8")

    # Recent Applications card
    rect = FancyBboxPatch((0.1, 2.65), 3.7, 2.1, boxstyle="round,pad=0.08",
                          linewidth=1, edgecolor="#E2E8F0", facecolor="white")
    ax.add_patch(rect)
    ax.text(0.3, 4.6, "Recent Applications", fontsize=9, fontweight="bold", color="#1E293B")
    apps = [("Google", "SWE", "#4F46E5", "Applied"),
            ("Stripe", "Backend", "#16A34A", "Offer"),
            ("Shopify", "Full Stack", "#7C3AED", "Interview")]
    for j, (co, role, col, st) in enumerate(apps):
        ry = 4.2 - j * 0.55
        ax.text(0.3, ry, co, fontsize=8, fontweight="bold", color="#334155")
        ax.text(0.3, ry - 0.22, role, fontsize=7, color="#94A3B8")
        badge = FancyBboxPatch((2.8, ry - 0.18), 0.88, 0.3,
                               boxstyle="round,pad=0.04", linewidth=0, facecolor=col+"22")
        ax.add_patch(badge)
        ax.text(3.24, ry - 0.02, st, ha="center", va="center",
                fontsize=6.5, color=col, fontweight="bold")

    # AI Activity card
    rect2 = FancyBboxPatch((3.95, 2.65), 1.9, 2.1, boxstyle="round,pad=0.08",
                           linewidth=1, edgecolor="#E2E8F0", facecolor="white")
    ax.add_patch(rect2)
    ax.text(4.1, 4.6, "AI Activity", fontsize=9, fontweight="bold", color="#1E293B")
    acts = ["Resume analyzed", "Job matched", "Cover letter\ngenerated", "Interview prep"]
    act_cols = ["#4F46E5", "#7C3AED", "#EC4899", "#D97706"]
    for j, (act, col) in enumerate(zip(acts, act_cols)):
        ay = 4.15 - j * 0.52
        dot2 = plt.Circle((4.18, ay), 0.07, color=col)
        ax.add_patch(dot2)
        ax.text(4.35, ay, act, fontsize=7, color="#475569", va="center")

    # Quick Actions
    rect3 = FancyBboxPatch((0.1, 0.15), 5.75, 2.35, boxstyle="round,pad=0.08",
                           linewidth=1, edgecolor="#E2E8F0", facecolor="white")
    ax.add_patch(rect3)
    ax.text(0.3, 2.38, "Quick Actions", fontsize=9, fontweight="bold", color="#1E293B")
    qas = [("Analyze Resume", "#4F46E5"), ("Add Application", "#2563EB"),
           ("Match to Job", "#7C3AED"), ("Cover Letter", "#059669")]
    for j, (label, col) in enumerate(qas):
        qx = 0.18 + j * 1.45
        rect4 = FancyBboxPatch((qx, 0.25), 1.32, 1.9, boxstyle="round,pad=0.06",
                               linewidth=0, facecolor=col+"18")
        ax.add_patch(rect4)
        circ = plt.Circle((qx + 0.66, 1.6), 0.3, color=col)
        ax.add_patch(circ)
        ax.text(qx + 0.66, 0.8, label, ha="center", fontsize=7,
                color=col, fontweight="bold", multialignment="center",
                wrap=True)

    add_fig(sl, fig, Inches(6.9), Inches(1.2), Inches(6.0), Inches(6.0))
    return sl


# ── Slide 6 · Resume Analyzer ─────────────────────────────────────────────
def slide_resume(prs):
    sl = blank_slide(prs)
    white_bg(sl)
    accent_strip(sl)
    slide_header(sl, "Feature 2 · Resume Analyzer",
                 "Upload once → get a full ATS audit powered by Claude")

    # Score circles diagram
    fig, axes = plt.subplots(1, 3, figsize=(5.5, 3.2))
    fig.patch.set_facecolor("white")
    scores = [(82, "ATS Score", "#4F46E5"), (74, "Formatting", "#7C3AED"), (88, "Content", "#059669")]
    for ax, (score, label, col) in zip(axes, scores):
        ax.set_facecolor("white")
        ax.set_xlim(-1.3, 1.3); ax.set_ylim(-1.3, 1.3); ax.set_aspect("equal"); ax.axis("off")
        theta = np.linspace(np.pi/2, np.pi/2 - 2*np.pi*(score/100), 100)
        btheta = np.linspace(0, 2*np.pi, 100)
        ax.plot(np.cos(btheta), np.sin(btheta), color="#F1F5F9", lw=10, solid_capstyle="round")
        ax.plot(np.cos(theta), np.sin(theta), color=col, lw=10, solid_capstyle="round")
        ax.text(0, 0.1, f"{score}%", ha="center", va="center",
                fontsize=18, fontweight="bold", color=col)
        ax.text(0, -0.38, label, ha="center", va="center",
                fontsize=9, color="#64748B")
    plt.tight_layout(pad=0.5)
    add_fig(sl, fig, Inches(6.8), Inches(1.25), Inches(6.0), Inches(3.4))

    # Output cards grid
    outputs = [
        ("#4F46E5", "Strengths", "Clear contact info, strong technical skills, quantified achievements"),
        ("#DC2626", "Weaknesses", "Missing summary, sparse keywords, formatting inconsistencies"),
        ("#D97706", "Recommendations", "High: Add professional summary  ·  Medium: Quantify more bullets"),
        ("#7C3AED", "Missing Keywords", "Docker · Kubernetes · CI/CD · Agile · REST APIs"),
    ]
    for i, (col, title, body) in enumerate(outputs):
        cx = Inches(6.8 + (i % 2) * 3.1)
        cy = Inches(4.55 + (i // 2) * 1.4)
        card = sl.shapes.add_shape(5, cx, cy, Inches(2.95), Inches(1.28))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)])
        card.line.fill.background()
        card.adjustments[0] = 0.06
        # Fake alpha by using light color
        card.fill.fore_color.rgb = SLATE_100
        card.line.color.rgb = RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)])
        card.line.width = Pt(1.5)
        add_text_box(sl, title, cx + Inches(0.12), cy + Inches(0.08), Inches(2.7), Inches(0.32),
                     font_size=11, bold=True,
                     color=RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)]))
        add_text_box(sl, body, cx + Inches(0.12), cy + Inches(0.38), Inches(2.7), Inches(0.82),
                     font_size=9.5, color=SLATE_700)

    # Left steps
    steps = [
        "Parse PDF / DOCX / TXT",
        "Extract experience & education",
        "Identify technical + soft skills",
        "Evaluate ATS compatibility",
        "Calculate 3 scores",
        "Generate prioritized recommendations",
    ]
    add_text_box(sl, "How It Works", Inches(0.45), Inches(1.45), Inches(6), Inches(0.4),
                 font_size=15, bold=True, color=SLATE_900)
    for i, step in enumerate(steps):
        y = Inches(1.95 + i * 0.76)
        dot = sl.shapes.add_shape(9, Inches(0.5), y + Inches(0.06),
                                  Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = INDIGO
        dot.line.fill.background()
        num_box = add_text_box(sl, str(i + 1), Inches(0.47), y + Inches(0.02),
                               Inches(0.22), Inches(0.26),
                               font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(sl, step, Inches(0.85), y, Inches(5.6), Inches(0.42),
                     font_size=13, color=SLATE_700)

    return sl


# ── Slide 7 · Job Tracker ─────────────────────────────────────────────────
def slide_tracker(prs):
    sl = blank_slide(prs)
    white_bg(sl)
    accent_strip(sl)
    slide_header(sl, "Feature 3 · Job Tracker",
                 "Full CRUD application — track every application in one place")

    # Left features
    feats = [
        ("Add / Edit / Delete", "Full modal form: company, role, status, location, URL, salary, notes, dates"),
        ("Status Filter Tabs", "All · Applied · Interviewing · Offer · Rejected · Saved with live counts"),
        ("Inline Status Update", "Click the status badge in the table row to change it instantly"),
        ("Bulk Operations", "Select multiple rows → delete in one action"),
        ("Demo Data Loader", "Seed realistic sample applications with one click"),
        ("Deep Links to AI", "Launch Job Matcher or Cover Letter Generator from any row"),
    ]
    for i, (title, desc) in enumerate(feats):
        y = Inches(1.45 + i * 0.95)
        dot = sl.shapes.add_shape(9, Inches(0.45), y + Inches(0.08),
                                  Inches(0.2), Inches(0.2))
        dot.fill.solid(); dot.fill.fore_color.rgb = INDIGO; dot.line.fill.background()
        add_text_box(sl, title, Inches(0.8), y, Inches(5.5), Inches(0.32),
                     font_size=13, bold=True, color=INDIGO)
        add_text_box(sl, desc, Inches(0.8), y + Inches(0.3), Inches(5.5), Inches(0.55),
                     font_size=11.5, color=SLATE_700)

    # Right: table mockup
    fig, ax = plt.subplots(figsize=(6.2, 5.5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.set_xlim(0, 7); ax.set_ylim(0, 6); ax.axis("off")

    # Filter tabs
    tabs = ["All 12", "Applied 5", "Offer 1", "Rejected 3", "Saved 2"]
    tab_cols = ["#4F46E5", "#94A3B8", "#94A3B8", "#94A3B8", "#94A3B8"]
    for j, (tab, tc) in enumerate(zip(tabs, tab_cols)):
        txc = "white" if j == 0 else "#64748B"
        bg  = "#4F46E5" if j == 0 else "#F1F5F9"
        rect = FancyBboxPatch((0.1 + j * 1.37, 5.4), 1.25, 0.45,
                              boxstyle="round,pad=0.05", linewidth=0, facecolor=bg)
        ax.add_patch(rect)
        ax.text(0.72 + j * 1.37, 5.63, tab, ha="center", va="center",
                fontsize=7.5, color=txc, fontweight="bold" if j == 0 else "normal")

    # Table header
    rect_h = FancyBboxPatch((0.1, 4.8), 6.8, 0.5, boxstyle="round,pad=0.04",
                            linewidth=0, facecolor="#F8FAFC")
    ax.add_patch(rect_h)
    for j, (hdr, hx) in enumerate(zip(["Company", "Role", "Status", "Date", "Actions"],
                                       [0.4, 1.8, 3.2, 4.5, 5.6])):
        ax.text(hx, 5.05, hdr, fontsize=7.5, color="#94A3B8", fontweight="bold")

    # Table rows
    rows = [
        ("Google", "Senior SWE", "#2563EB", "Applied",  "Jan 15"),
        ("Stripe",  "Backend Eng","#16A34A", "Offer",    "Jan 10"),
        ("Shopify", "Full Stack",  "#7C3AED","Interview","Jan 08"),
        ("Netflix", "Data Eng",    "#DC2626","Rejected", "Dec 30"),
        ("Apple",   "iOS Dev",     "#64748B","Saved",    "—"),
    ]
    for j, (co, role, col, st, date) in enumerate(rows):
        ry = 4.55 - j * 0.78
        if j % 2 == 0:
            row_bg = FancyBboxPatch((0.1, ry - 0.38), 6.8, 0.72,
                                   boxstyle="round,pad=0.02", linewidth=0, facecolor="#FAFAFA")
            ax.add_patch(row_bg)
        ax.text(0.4, ry, co,    fontsize=8.5, color="#1E293B", fontweight="bold", va="center")
        ax.text(1.8, ry, role,  fontsize=8,   color="#475569", va="center")
        badge = FancyBboxPatch((3.1, ry - 0.17), 1.05, 0.34,
                               boxstyle="round,pad=0.04", linewidth=0, facecolor=col+"22")
        ax.add_patch(badge)
        ax.text(3.62, ry, st, ha="center", va="center",
                fontsize=7, color=col, fontweight="bold")
        ax.text(4.5, ry, date, fontsize=7.5, color="#94A3B8", va="center")
        # action icons text
        ax.text(5.6, ry, "⊕ ✎ ✕", fontsize=8, color="#94A3B8", va="center")

    add_fig(sl, fig, Inches(6.8), Inches(1.2), Inches(6.1), Inches(6.0))
    return sl


# ── Slide 8 · Job Matcher ─────────────────────────────────────────────────
def slide_matcher(prs):
    sl = blank_slide(prs)
    white_bg(sl)
    accent_strip(sl)
    slide_header(sl, "Feature 4 · Job Matcher",
                 "Paste resume + job description → Claude scores the fit and tells you exactly what to fix")

    # Left: what you get
    sections = [
        ("Match Score & Likelihood", "#4F46E5",
         "0–100% score with Low / Medium / High likelihood badge"),
        ("Matching vs Missing Skills", "#16A34A",
         "Color-coded skill tags — green matched, red missing"),
        ("Experience Gaps", "#D97706",
         "Side-by-side matching experience and gap analysis"),
        ("Tailoring Suggestions", "#7C3AED",
         "Section-by-section resume advice tied to the JD"),
        ("Better Bullets  (AI)", "#0891B2",
         "Rewrites your resume bullets to match the JD keywords"),
        ("Interview Prep  (AI)", "#EA580C",
         "Behavioral/technical/situational Q&A accordion"),
    ]
    for i, (title, col, desc) in enumerate(sections):
        y = Inches(1.45 + i * 0.95)
        tag = sl.shapes.add_shape(5, Inches(0.42), y, Inches(0.16), Inches(0.36))
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)])
        tag.line.fill.background()
        add_text_box(sl, title, Inches(0.72), y, Inches(5.4), Inches(0.33),
                     font_size=13, bold=True,
                     color=RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)]))
        add_text_box(sl, desc, Inches(0.72), y + Inches(0.31), Inches(5.4), Inches(0.55),
                     font_size=11.5, color=SLATE_700)

    # Right: match result mockup
    fig, ax = plt.subplots(figsize=(5.8, 5.8))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.set_xlim(0, 6); ax.set_ylim(0, 6.2); ax.axis("off")

    # Big score circle
    theta = np.linspace(np.pi/2, np.pi/2 - 2*np.pi*0.78, 100)
    btheta = np.linspace(0, 2*np.pi, 100)
    ax.plot(np.cos(btheta) * 0.9 + 1.3, np.sin(btheta) * 0.9 + 5.0,
            color="#F1F5F9", lw=14, solid_capstyle="round")
    ax.plot(np.cos(theta) * 0.9 + 1.3, np.sin(theta) * 0.9 + 5.0,
            color="#4F46E5", lw=14, solid_capstyle="round")
    ax.text(1.3, 5.1, "78%", ha="center", va="center",
            fontsize=22, fontweight="bold", color="#4F46E5")
    ax.text(1.3, 4.55, "Match Score", ha="center", va="center",
            fontsize=9, color="#64748B")

    # Likelihood badge
    badge = FancyBboxPatch((2.3, 4.7), 1.4, 0.45, boxstyle="round,pad=0.06",
                           linewidth=0, facecolor="#D1FAE5")
    ax.add_patch(badge)
    ax.text(3.0, 4.93, "High likelihood", ha="center", va="center",
            fontsize=9, color="#16A34A", fontweight="bold")

    ax.text(2.3, 4.3, "Strong technical alignment. Add Docker and\nCI/CD keywords to push above 85%.",
            fontsize=8, color="#475569", va="top")

    # Skills grid
    matching = ["Python", "FastAPI", "React", "PostgreSQL", "REST API"]
    missing  = ["Docker", "Kubernetes", "CI/CD", "Redis"]
    ax.text(0.1, 3.9, "Matching Skills", fontsize=9, fontweight="bold", color="#16A34A")
    for j, sk in enumerate(matching):
        bx = 0.1 + (j % 3) * 1.2
        by = 3.5 - (j // 3) * 0.45
        b = FancyBboxPatch((bx, by), 1.05, 0.32, boxstyle="round,pad=0.05",
                           linewidth=0, facecolor="#DCFCE7")
        ax.add_patch(b)
        ax.text(bx + 0.52, by + 0.16, sk, ha="center", va="center",
                fontsize=7.5, color="#16A34A", fontweight="bold")

    ax.text(0.1, 2.8, "Missing Skills", fontsize=9, fontweight="bold", color="#DC2626")
    for j, sk in enumerate(missing):
        bx = 0.1 + j * 1.4
        b2 = FancyBboxPatch((bx, 2.4), 1.25, 0.32, boxstyle="round,pad=0.05",
                            linewidth=0, facecolor="#FEE2E2")
        ax.add_patch(b2)
        ax.text(bx + 0.62, 2.56, sk, ha="center", va="center",
                fontsize=7.5, color="#DC2626", fontweight="bold")

    # Action buttons
    btns = [("Save to Tracker","#4F46E5"),("Better Bullets","#7C3AED"),
            ("Cover Letter","#059669"),("Interview Prep","#D97706")]
    for j, (lbl, col) in enumerate(btns):
        bx = 0.1 + (j % 2) * 3.0
        by = 1.55 - (j // 2) * 0.62
        b = FancyBboxPatch((bx, by), 2.7, 0.48, boxstyle="round,pad=0.06",
                           linewidth=0, facecolor=col)
        ax.add_patch(b)
        ax.text(bx + 1.35, by + 0.24, lbl, ha="center", va="center",
                fontsize=8.5, color="white", fontweight="bold")

    add_fig(sl, fig, Inches(6.9), Inches(1.15), Inches(6.0), Inches(6.1))
    return sl


# ── Slide 9 · Cover Letter Generator ─────────────────────────────────────
def slide_cover_letter(prs):
    sl = blank_slide(prs)
    white_bg(sl)
    accent_strip(sl)
    slide_header(sl, "Feature 5 · Cover Letter Generator",
                 "AI-written, perfectly tailored — copy, download, or save to your application")

    # Tone selector diagram
    fig, ax = plt.subplots(figsize=(5.8, 2.2))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.set_xlim(0, 6); ax.set_ylim(0, 2.4); ax.axis("off")
    tones = [("Professional","#4F46E5",True), ("Enthusiastic","#2563EB",False),
             ("Concise","#7C3AED",False), ("Creative","#EC4899",False), ("Storytelling","#D97706",False)]
    for j, (tone, col, selected) in enumerate(tones):
        bx = 0.05 + j * 1.18
        bg = col if selected else "white"
        tc = "white" if selected else "#64748B"
        ec = col
        b = FancyBboxPatch((bx, 0.8), 1.1, 0.75, boxstyle="round,pad=0.06",
                           linewidth=2, edgecolor=ec, facecolor=bg)
        ax.add_patch(b)
        ax.text(bx + 0.55, 1.17, tone, ha="center", va="center",
                fontsize=8.5, color=tc, fontweight="bold" if selected else "normal")
    ax.text(3.0, 2.2, "Tone Selector", ha="center", va="top",
            fontsize=10, fontweight="bold", color="#1E293B")
    ax.text(3.0, 0.45, "← Select Professional for corporate roles · Creative for startups →",
            ha="center", va="center", fontsize=7.5, color="#94A3B8", style="italic")
    add_fig(sl, fig, Inches(6.8), Inches(1.2), Inches(6.0), Inches(2.4))

    # Letter preview mockup
    fig2, ax2 = plt.subplots(figsize=(5.8, 3.5))
    fig2.patch.set_facecolor("white")
    ax2.set_facecolor("white")
    ax2.set_xlim(0, 6); ax2.set_ylim(0, 3.8); ax2.axis("off")
    rect = FancyBboxPatch((0.1, 0.2), 5.8, 3.45, boxstyle="round,pad=0.08",
                          linewidth=1, edgecolor="#E2E8F0", facecolor="white")
    ax2.add_patch(rect)
    ax2.text(3.0, 3.5, "Cover Letter  —  Google  ·  Professional  ·  342 words",
             ha="center", va="top", fontsize=8, color="#64748B")
    letter_lines = [
        "Dear Hiring Manager,",
        "",
        "I am excited to apply for the Senior Software Engineer",
        "position at Google. With 5 years of experience building",
        "scalable systems with Python and React, I am confident...",
        "",
        "In my current role at Acme Corp I led the migration of",
        "a monolithic service to microservices, reducing latency",
        "by 40% and improving deployment frequency 3×...",
        "",
        "I look forward to the opportunity to contribute to Google's",
        "mission and bring my expertise to your team.",
        "",
        "Sincerely,  Jane Smith",
    ]
    for j, line in enumerate(letter_lines):
        ax2.text(0.35, 3.15 - j * 0.215, line, fontsize=7.5, color="#334155", va="top")

    # Action buttons inside figure
    for j, (lbl, col) in enumerate([("Copy","#4F46E5"),("Download","#64748B"),("Regenerate","#94A3B8")]):
        bx = 0.2 + j * 1.95
        b = FancyBboxPatch((bx, 0.0), 1.75, 0.38, boxstyle="round,pad=0.05",
                           linewidth=0, facecolor=col)
        ax2.add_patch(b)
        ax2.text(bx+0.87, 0.19, lbl, ha="center", va="center",
                 fontsize=8.5, color="white", fontweight="bold")

    add_fig(sl, fig2, Inches(6.8), Inches(3.55), Inches(6.0), Inches(3.7))

    # Left features
    features = [
        ("5 Tone Options", "Professional · Enthusiastic · Concise · Creative · Storytelling"),
        ("Context Pre-fill", "Resume and JD auto-loaded from localStorage — no re-typing"),
        ("Job Selector", "Pick any saved application from the tracker to pre-fill details"),
        ("Copy / Download", "Clipboard copy or .txt download with one click"),
        ("Regenerate", "Instant re-generation with a different tone via shortcut buttons"),
        ("Save to Application", "Appends letter snippet to the job's notes in the tracker"),
    ]
    for i, (title, desc) in enumerate(features):
        y = Inches(1.45 + i * 0.97)
        tag = sl.shapes.add_shape(9, Inches(0.45), y + Inches(0.07),
                                  Inches(0.18), Inches(0.18))
        tag.fill.solid(); tag.fill.fore_color.rgb = INDIGO; tag.line.fill.background()
        add_text_box(sl, title, Inches(0.8), y, Inches(5.5), Inches(0.33),
                     font_size=13, bold=True, color=INDIGO)
        add_text_box(sl, desc, Inches(0.8), y + Inches(0.31), Inches(5.5), Inches(0.56),
                     font_size=11.5, color=SLATE_700)

    return sl


# ── Slide 10 · AI Agent UX ────────────────────────────────────────────────
def slide_agent_ux(prs):
    sl = blank_slide(prs)
    dark_bg(sl)
    add_rect(sl, 0, 0, Inches(0.18), H, INDIGO)
    slide_header(sl, "AI Agent UX Pattern",
                 "Every Claude call surfaces its work — not just a spinner", dark=True)

    # Left: explanation
    points = [
        ("Transparent Processing",
         "Users see exactly what Claude is doing at each step — parsing, scoring, generating."),
        ("Reusable Component",
         "AgentActivity accepts a steps[] array, isRunning, and isDone flags — used across all 4 AI flows."),
        ("Trust Through Visibility",
         "Showing step-by-step progress builds confidence that the AI is working, not stalled."),
        ("Animated Completion",
         "Each step gets a checkmark as it completes. The panel stays visible after finishing so results feel earned."),
    ]
    for i, (title, desc) in enumerate(points):
        y = Inches(1.55 + i * 1.35)
        dot = sl.shapes.add_shape(9, Inches(0.5), y + Inches(0.09),
                                  Inches(0.2), Inches(0.2))
        dot.fill.solid(); dot.fill.fore_color.rgb = INDIGO; dot.line.fill.background()
        add_text_box(sl, title, Inches(0.85), y, Inches(5.4), Inches(0.36),
                     font_size=14, bold=True, color=WHITE)
        add_text_box(sl, desc, Inches(0.85), y + Inches(0.37), Inches(5.4), Inches(0.85),
                     font_size=12, color=SLATE_500)

    # Right: agent activity UI mockup
    fig, ax = plt.subplots(figsize=(5.8, 5.6))
    fig.patch.set_facecolor("#1E293B")
    ax.set_facecolor("#1E293B")
    ax.set_xlim(0, 6); ax.set_ylim(0, 6); ax.axis("off")

    rect_card = FancyBboxPatch((0.15, 0.4), 5.7, 5.3, boxstyle="round,pad=0.1",
                               linewidth=1, edgecolor="#334155", facecolor="#0F172A")
    ax.add_patch(rect_card)

    ax.text(3.0, 5.55, "Resume Analyzer", ha="center", va="top",
            fontsize=11, fontweight="bold", color="#A5B4FC")
    ax.text(3.0, 5.18, "Claude is analyzing your resume…",
            ha="center", va="top", fontsize=8.5, color="#64748B")

    steps_data = [
        ("Parsing resume document",         True),
        ("Extracting work experience",       True),
        ("Identifying technical skills",     True),
        ("Evaluating ATS compatibility",     True),
        ("Calculating scores",               False),  # animating
        ("Generating recommendations",       False),
    ]
    for j, (step, done) in enumerate(steps_data):
        sy = 4.65 - j * 0.65
        # Icon
        if done:
            circ = plt.Circle((0.6, sy + 0.12), 0.18, color="#4F46E5")
            ax.add_patch(circ)
            ax.text(0.6, sy + 0.12, "✓", ha="center", va="center",
                    fontsize=8, color="white", fontweight="bold")
        else:
            circ2 = plt.Circle((0.6, sy + 0.12), 0.18, color="#334155")
            ax.add_patch(circ2)
            if j == 4:  # animating dot
                for k in range(3):
                    d = plt.Circle((0.45 + k*0.12, sy + 0.12), 0.04, color="#64748B")
                    ax.add_patch(d)

        col = "#E2E8F0" if done else "#64748B"
        ax.text(0.92, sy + 0.12, step, va="center", fontsize=9,
                color=col, fontweight="bold" if done else "normal")

    # Progress bar
    ax.add_patch(FancyBboxPatch((0.35, 0.55), 5.3, 0.25,
                                boxstyle="round,pad=0.04", linewidth=0, facecolor="#1E293B"))
    ax.add_patch(FancyBboxPatch((0.35, 0.55), 5.3 * 0.67, 0.25,
                                boxstyle="round,pad=0.04", linewidth=0, facecolor="#4F46E5"))
    ax.text(5.65, 0.67, "67%", va="center", ha="right", fontsize=8,
            color="#A5B4FC", fontweight="bold")

    add_fig(sl, fig, Inches(6.9), Inches(1.15), Inches(6.0), Inches(6.1))
    return sl


# ── Slide 11 · Engineering Decisions ──────────────────────────────────────
def slide_engineering(prs):
    sl = blank_slide(prs)
    white_bg(sl)
    accent_strip(sl)
    slide_header(sl, "Key Engineering Decisions",
                 "The non-obvious choices that make it work smoothly")

    decisions = [
        ("#4F46E5", "localStorage for Context Passing",
         "Resume text and JD are cached in the browser so navigating from Analyzer → "
         "Matcher → Cover Letter is seamless — no re-uploading or re-pasting required."),
        ("#7C3AED", "Structured JSON from Claude",
         "All Claude responses are requested as typed JSON schemas so the frontend can "
         "render individual sections (scores, skill tags, recommendations) instead of raw text."),
        ("#16A34A", "Single claude_service.py",
         "All Anthropic API calls live in one service file, making it easy to swap "
         "models, tune prompts, or add caching in a single place."),
        ("#D97706", "Next.js API Proxy (catch-all route)",
         "A catch-all App Router route proxies /api/* to the FastAPI backend, eliminating "
         "CORS issues in production and keeping env vars server-side only."),
        ("#2563EB", "SQLite → PostgreSQL in one line",
         "SQLAlchemy abstracts the database; swapping to Postgres for production "
         "requires changing only the connection string."),
        ("#EC4899", "Composable AI Flows",
         "The Job Matcher chains 4 distinct Claude calls (match → bullets → interview Q&A "
         "→ cover letter navigation) from a single result page without overwhelming the user."),
    ]

    for i, (col, title, desc) in enumerate(decisions):
        row, c = divmod(i, 2)
        x = Inches(0.4 + c * 6.4)
        y = Inches(1.45 + row * 2.05)

        card = sl.shapes.add_shape(5, x, y, Inches(6.1), Inches(1.9))
        card.fill.solid(); card.fill.fore_color.rgb = SLATE_100
        card.line.color.rgb = RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)])
        card.line.width = Pt(2); card.adjustments[0] = 0.06

        # Left accent bar
        bar = sl.shapes.add_shape(1, x, y, Inches(0.07), Inches(1.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)])
        bar.line.fill.background()

        add_text_box(sl, title, x + Inches(0.18), y + Inches(0.1), Inches(5.8), Inches(0.36),
                     font_size=13, bold=True,
                     color=RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)]))
        add_text_box(sl, desc, x + Inches(0.18), y + Inches(0.46), Inches(5.8), Inches(1.3),
                     font_size=11.5, color=SLATE_700, wrap=True)

    return sl


# ── Slide 12 · Demo Flow ──────────────────────────────────────────────────
def slide_demo(prs):
    sl = blank_slide(prs)
    dark_bg(sl)
    add_rect(sl, 0, 0, Inches(0.18), H, INDIGO)
    slide_header(sl, "Live Demo Walk-through",
                 "6 steps that showcase the full user journey", dark=True)

    steps = [
        ("#4F46E5", "1. Dashboard",
         "Open the app — see empty state with stat cards and quick action shortcuts."),
        ("#2563EB", "2. Resume Analyzer",
         "Upload a PDF → watch agent steps animate → get ATS score, skills, recommendations."),
        ("#7C3AED", "3. Job Matcher",
         "Resume pre-filled → paste any JD → analyze → generate bullets & interview Q&A."),
        ("#059669", "4. Save to Tracker",
         "Hit 'Save to Tracker' → fill company/role → job saved with AI match score in notes."),
        ("#D97706", "5. Job Tracker",
         "See the new entry → change status inline → use filter tabs to slice the list."),
        ("#EC4899", "6. Cover Letter",
         "Click Cover Letter icon on the row → context pre-filled → choose Creative → generate → copy."),
    ]

    for i, (col, title, desc) in enumerate(steps):
        row, c = divmod(i, 2)
        x = Inches(0.4 + c * 6.4)
        y = Inches(1.45 + row * 2.0)

        card = sl.shapes.add_shape(5, x, y, Inches(6.1), Inches(1.78))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(30, 41, 59)   # slate-800
        card.line.color.rgb = RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)])
        card.line.width = Pt(2); card.adjustments[0] = 0.06

        bar = sl.shapes.add_shape(1, x, y, Inches(0.07), Inches(1.78))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)])
        bar.line.fill.background()

        add_text_box(sl, title, x + Inches(0.18), y + Inches(0.1), Inches(5.8), Inches(0.36),
                     font_size=14, bold=True,
                     color=RGBColor(*[int(col[i:i+2], 16) for i in (1, 3, 5)]))
        add_text_box(sl, desc, x + Inches(0.18), y + Inches(0.46), Inches(5.8), Inches(1.1),
                     font_size=12, color=SLATE_500, wrap=True)

    return sl


# ── Slide 13 · Deployment ─────────────────────────────────────────────────
def slide_deployment(prs):
    sl = blank_slide(prs)
    white_bg(sl)
    accent_strip(sl)
    slide_header(sl, "Deployment", "Production-ready, zero-cost hosting on free tiers")

    fig, ax = plt.subplots(figsize=(12, 5.0))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.set_xlim(0, 13); ax.set_ylim(0, 5.5); ax.axis("off")

    # Three deployment boxes
    boxes = [
        (0.3, 1.2, 3.8, 3.6, "#EEF2FF", "#4F46E5", "Vercel", "Frontend",
         ["Next.js App Router", "Auto-deploy on push", "Edge CDN", "Free tier"]),
        (4.8, 1.2, 3.8, 3.6, "#F5F3FF", "#7C3AED", "Render", "Backend",
         ["FastAPI + Uvicorn", "render.yaml config", "Persistent disk", "Free tier"]),
        (9.3, 1.2, 3.4, 3.6, "#ECFDF5", "#059669", "Anthropic", "Claude API",
         ["claude-sonnet-4-6", "Pay-per-use", "JSON mode", "Streaming capable"]),
    ]
    for lx, ly, lw, lh, bg, col, name, sub, items in boxes:
        rect = FancyBboxPatch((lx, ly), lw, lh, boxstyle="round,pad=0.15",
                              linewidth=2.5, edgecolor=col, facecolor=bg)
        ax.add_patch(rect)
        ax.text(lx + lw/2, ly + lh - 0.25, name, ha="center", va="top",
                fontsize=15, fontweight="bold", color=col)
        ax.text(lx + lw/2, ly + lh - 0.72, sub, ha="center", va="top",
                fontsize=10, color="#64748B")
        for j, itm in enumerate(items):
            ax.text(lx + 0.25, ly + lh - 1.35 - j * 0.52, f"• {itm}",
                    fontsize=9.5, color="#334155", va="top")

    # Arrows
    ax.annotate("", xy=(4.8, 3.0), xytext=(4.1, 3.0),
                arrowprops=dict(arrowstyle="->", color="#94A3B8", lw=2))
    ax.text(4.45, 3.2, "REST\n/api/*", ha="center", fontsize=8, color="#94A3B8")

    ax.annotate("", xy=(9.3, 3.0), xytext=(8.6, 3.0),
                arrowprops=dict(arrowstyle="->", color="#94A3B8", lw=2))
    ax.text(8.95, 3.2, "Claude\nSDK", ha="center", fontsize=8, color="#94A3B8")

    # Env vars
    ev_box = FancyBboxPatch((0.3, 0.1), 12.4, 0.9, boxstyle="round,pad=0.1",
                            linewidth=1, edgecolor="#E2E8F0", facecolor="#F8FAFC")
    ax.add_patch(ev_box)
    ax.text(6.5, 0.92, "Environment Variables", ha="center", va="top",
            fontsize=10, fontweight="bold", color="#334155")
    ax.text(2.0, 0.55, "ANTHROPIC_API_KEY  (backend only)",
            fontsize=9, color="#64748B", ha="center")
    ax.text(6.5, 0.55, "NEXT_PUBLIC_API_URL  (→ Render URL)",
            fontsize=9, color="#64748B", ha="center")
    ax.text(11.0, 0.55, "DATABASE_URL  (SQLite / Postgres)",
            fontsize=9, color="#64748B", ha="center")

    add_fig(sl, fig, Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.8))
    return sl


# ── Slide 14 · Thank You ──────────────────────────────────────────────────
def slide_thank_you(prs):
    sl = blank_slide(prs)
    dark_bg(sl)
    add_rect(sl, 0, 0, Inches(0.18), H, INDIGO)

    add_text_box(sl, "Thank You",
                 Inches(0.6), Inches(1.5), Inches(9), Inches(1.4),
                 font_size=56, bold=True, color=WHITE)
    add_text_box(sl, "AI Job Analyzer",
                 Inches(0.6), Inches(2.95), Inches(9), Inches(0.6),
                 font_size=24, color=RGBColor(0xA5, 0xB4, 0xFC), bold=True)
    add_text_box(sl,
                 "Built with Next.js 14  ·  FastAPI  ·  Claude claude-sonnet-4-6  ·  Tailwind CSS",
                 Inches(0.6), Inches(3.65), Inches(10), Inches(0.45),
                 font_size=15, color=SLATE_500)

    # Summary bullets
    summary = [
        "Resume Analyzer  — ATS scores, skill extraction, prioritized recommendations",
        "Job Tracker       — Full CRUD, inline status updates, bulk operations",
        "Job Matcher       — Fit score, gap analysis, AI bullets, interview prep",
        "Cover Letter Gen  — 5 tones, pre-filled context, copy / download / save",
    ]
    for i, line in enumerate(summary):
        y = Inches(4.35 + i * 0.6)
        dot = sl.shapes.add_shape(9, Inches(0.6), y + Inches(0.1),
                                  Inches(0.16), Inches(0.16))
        dot.fill.solid(); dot.fill.fore_color.rgb = INDIGO; dot.line.fill.background()
        add_text_box(sl, line, Inches(0.92), y, Inches(10), Inches(0.42),
                     font_size=13, color=SLATE_500)

    add_text_box(sl, "Questions?",
                 Inches(9.5), Inches(6.5), Inches(3.5), Inches(0.6),
                 font_size=22, bold=True, color=INDIGO, align=PP_ALIGN.RIGHT)

    return sl


# ═══════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    print("Building slides…")
    slide_title(prs);       print("  1/14  Title")
    slide_problem(prs);     print("  2/14  The Problem")
    slide_tech_stack(prs);  print("  3/14  Tech Stack")
    slide_architecture(prs);print("  4/14  Architecture")
    slide_dashboard(prs);   print("  5/14  Dashboard")
    slide_resume(prs);      print("  6/14  Resume Analyzer")
    slide_tracker(prs);     print("  7/14  Job Tracker")
    slide_matcher(prs);     print("  8/14  Job Matcher")
    slide_cover_letter(prs);print("  9/14  Cover Letter")
    slide_agent_ux(prs);    print(" 10/14  AI Agent UX")
    slide_engineering(prs); print(" 11/14  Engineering Decisions")
    slide_demo(prs);         print(" 12/14  Demo Flow")
    slide_deployment(prs);  print(" 13/14  Deployment")
    slide_thank_you(prs);   print(" 14/14  Thank You")

    out = r"e:\projects\AI Projects\AIJob-Analyzer\AIJobAnalyzer.pptx"
    prs.save(out)
    print(f"\nSaved -> {out}")


if __name__ == "__main__":
    build()
